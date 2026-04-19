"""Extract text from office documents and PDFs."""
from __future__ import annotations

import io
import logging
from typing import Iterable

log = logging.getLogger(__name__)


def _from_pdf(data: bytes) -> str:
    import pdfplumber
    out: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            txt = page.extract_text() or ""
            if txt.strip():
                out.append(f"--- หน้า {i} ---\n{txt.strip()}")
            for tbl in page.extract_tables() or []:
                rows = ["\t".join((c or "") for c in row) for row in tbl]
                out.append("\n".join(rows))
    return "\n\n".join(out)


def _from_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for tbl in doc.tables:
        for row in tbl.rows:
            parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def _from_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"--- ชีต: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append("\t".join(cells))
    return "\n".join(out)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"--- สไลด์ {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip()
                    if txt:
                        out.append(txt)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                out.append(f"[หมายเหตุ] {note}")
    return "\n".join(out)


_EXTRACTORS = {
    "pdf":  _from_pdf,
    "docx": _from_docx,
    "xlsx": _from_xlsx,
    "pptx": _from_pptx,
}

SUPPORTED_EXTS: Iterable[str] = tuple(_EXTRACTORS.keys())


def extract(data: bytes, ext: str) -> str:
    """Return extracted text for supported document types; '' otherwise."""
    fn = _EXTRACTORS.get(ext.lower())
    if not fn:
        return ""
    try:
        return fn(data) or ""
    except Exception as e:
        log.warning("doc extract failed (%s): %s", ext, e)
        return ""
